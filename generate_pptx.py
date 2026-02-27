"""
AquaGuard — PowerPoint Presentation Generator
Generates a professional 7-slide deck summarizing the solution.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(__file__), "docs", "AquaGuard_Presentation.pptx")

# ── Color palette ───────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
BG_CARD   = RGBColor(0x16, 0x21, 0x3E)   # card bg
ACCENT    = RGBColor(0x4E, 0xCD, 0xC4)   # teal
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xBB, 0xC5, 0xD5)
RED       = RGBColor(0xFF, 0x00, 0x54)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
GREEN     = RGBColor(0x00, 0xD2, 0x6A)
BLUE      = RGBColor(0x00, 0x9D, 0xFF)


def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper to add a styled textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT, bullet_color=ACCENT):
    """Add a bullet-point list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Decorative accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(slide, 1.5, 1.5, 10, 1.2,
                "AquaGuard", font_size=54, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 2.7, 10, 0.8,
                "AI-Powered Real-Time Urban Water Quality Monitoring",
                font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.5), Inches(3.8), Inches(4.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide, 1.5, 4.2, 10, 0.6,
                "Degradation Prediction  |  Early Warning  |  Probable Cause Identification",
                font_size=18, color=LIGHT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 5.5, 10, 0.5,
                "AI Night Challenge 2026",
                font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════
    # SLIDE 2: The Problem
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 6, 0.7,
                "The Problem", font_size=36, color=ACCENT, bold=True)

    # Stats cards
    stats = [
        ("2B+", "people lack safe\ndrinking water"),
        ("485K", "deaths/year from\ncontaminated water"),
        ("Hours", "delay in current\nlab-based testing"),
        (">60%", "false alarm rate\nin threshold systems"),
    ]
    for i, (big, small) in enumerate(stats):
        x = 0.8 + i * 3.1
        add_rounded_rect(slide, x, 1.5, 2.8, 1.8)
        add_textbox(slide, x + 0.1, 1.6, 2.6, 0.8,
                    big, font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.3, 2.6, 0.8,
                    small, font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, 0.8, 3.8, 11, 3.5, [
        "Current monitoring: periodic grab samples, lab results in hours to days",
        "Threshold-based SCADA alarms: too many false positives, no prediction capability",
        "No automated root-cause identification -- operators rely on experience",
        "Gap: no early warning BEFORE water quality degrades to unsafe levels",
        "Public health at risk: contamination events detected AFTER consumption",
    ], font_size=16)

    # ═══════════════════════════════════════════════════
    # SLIDE 3: Our Solution — Architecture
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 8, 0.7,
                "Our Solution: End-to-End Architecture", font_size=36, color=ACCENT, bold=True)

    # Architecture flow — simplified boxes
    boxes = [
        (0.5, 1.8, "Multi-Sensor\nIngestion", "pH, Turbidity, Chloramine\nConductivity, THM, TDS\n5-min intervals, 5 sites"),
        (3.5, 1.8, "Feature\nEngineering", "279 features:\nRolling stats, EWMA\nDerivatives, Lags\nCyclical time encoding"),
        (6.5, 1.8, "ML Models\n(LightGBM)", "Potability classifier\nCause identifier\nPer-sensor forecasters"),
        (9.5, 1.8, "Risk Index\n& Alerts", "WQRI 0-100 score\nPersistence + Hysteresis\nFalse alarm reduction"),
    ]
    for x, y, title, desc in boxes:
        add_rounded_rect(slide, x, y, 2.7, 2.5)
        add_textbox(slide, x + 0.1, y + 0.1, 2.5, 0.7,
                    title, font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, y + 0.9, 2.5, 1.5,
                    desc, font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrows between boxes
    for i in range(3):
        x = 3.2 + i * 3.0
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        Inches(x), Inches(2.8), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

    # Bottom section: Dashboard + API
    add_rounded_rect(slide, 0.5, 4.8, 5.8, 1.8)
    add_textbox(slide, 0.6, 4.9, 5.6, 0.5,
                "Streamlit Dashboard (5 Tabs)", font_size=18, color=ACCENT, bold=True)
    add_bullet_list(slide, 0.8, 5.4, 5.4, 1.2, [
        "Risk Overview: gauge + timeline + multi-site comparison",
        "Sensor Trends: per-sensor charts with regulatory limits",
        "Alerts Log: transition history with false alarm tracking",
        "Root Cause: diagnosis panel with confidence + explanations",
    ], font_size=12)

    add_rounded_rect(slide, 6.8, 4.8, 5.8, 1.8)
    add_textbox(slide, 6.9, 4.9, 5.6, 0.5,
                "FastAPI Inference Service", font_size=18, color=ACCENT, bold=True)
    add_bullet_list(slide, 7.0, 5.4, 5.4, 1.2, [
        "POST /predict -- single reading prediction",
        "WS /ws/stream -- WebSocket live demo streaming",
        "GET /health -- system health check",
        "Production-ready: async, CORS, schema validation",
    ], font_size=12)

    # ═══════════════════════════════════════════════════
    # SLIDE 4: Risk Index & Detection
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 10, 0.7,
                "Water Quality Risk Index (WQRI)", font_size=36, color=ACCENT, bold=True)

    # Formula
    add_rounded_rect(slide, 0.5, 1.3, 12, 1.0)
    add_textbox(slide, 0.8, 1.4, 11.5, 0.8,
                "WQRI = EWMA( SUM(weight_i x deviation_score_i) / SUM(weight_i) )     score: 0-100",
                font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Weights table
    add_textbox(slide, 0.5, 2.6, 5, 0.5,
                "Sensor Weights (domain-driven)", font_size=20, color=WHITE, bold=True)

    weights = [
        ("Chloramine", "0.22", "Disinfectant loss = pathogen risk"),
        ("Turbidity",  "0.20", "Pathogen surrogate"),
        ("THM",        "0.15", "Carcinogenic byproduct"),
        ("pH",         "0.12", "Corrosion / scaling"),
        ("Conductivity","0.10", "Contamination proxy"),
        ("Organic C",  "0.08", "DBP precursor"),
        ("TDS / Sulfate / Hardness", "0.13", "Secondary quality"),
    ]
    for i, (sensor, weight, reason) in enumerate(weights):
        y = 3.2 + i * 0.38
        add_textbox(slide, 0.8, y, 2.0, 0.35, sensor, font_size=13, color=WHITE)
        add_textbox(slide, 2.9, y, 0.8, 0.35, weight, font_size=13, color=ACCENT, bold=True)
        add_textbox(slide, 3.8, y, 4.0, 0.35, reason, font_size=12, color=LIGHT)

    # Risk scale
    add_textbox(slide, 7.5, 2.6, 5, 0.5,
                "Risk Scale & Alert Logic", font_size=20, color=WHITE, bold=True)

    scales = [
        ("0-30",  "EXCELLENT", GREEN,  "Routine monitoring"),
        ("30-50", "NORMAL",    BLUE,   "Minor deviations"),
        ("50-65", "CAUTION",   RGBColor(0xFF, 0xC1, 0x07), "Increased monitoring"),
        ("65-80", "WARNING",   ORANGE, "Investigate immediately"),
        ("80-100","CRITICAL",  RED,    "Emergency response"),
    ]
    for i, (rng, label, clr, action) in enumerate(scales):
        y = 3.2 + i * 0.45
        add_textbox(slide, 7.8, y, 1.0, 0.4, rng, font_size=14, color=clr, bold=True)
        add_textbox(slide, 8.9, y, 1.5, 0.4, label, font_size=13, color=clr, bold=True)
        add_textbox(slide, 10.5, y, 2.3, 0.4, action, font_size=12, color=LIGHT)

    # False alarm reduction
    add_rounded_rect(slide, 7.5, 5.6, 5.2, 1.3)
    add_textbox(slide, 7.7, 5.7, 5.0, 0.4,
                "False Alarm Reduction", font_size=16, color=ACCENT, bold=True)
    add_bullet_list(slide, 7.7, 6.1, 5.0, 0.8, [
        "Persistence: 3 consecutive readings (15 min) required",
        "Hysteresis: 5-point buffer before alert clears",
        "Result: ~80% fewer false alarms vs raw thresholds",
    ], font_size=12)

    # ═══════════════════════════════════════════════════
    # SLIDE 5: Root Cause Identification
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 10, 0.7,
                "Probable Cause Identification", font_size=36, color=ACCENT, bold=True)

    add_textbox(slide, 0.8, 1.2, 10, 0.5,
                "Hybrid approach: ML Classifier + Domain Rule Engine + SHAP Explanations",
                font_size=18, color=LIGHT)

    # Cause table
    causes = [
        ("Disinfectant Decay",     "Chloramine drops, THM rises",            "Hours-days"),
        ("Contamination/Intrusion","Turbidity spike, conductivity jump",     "Minutes-hours"),
        ("Pipe Corrosion",         "pH drops, hardness rises, conductivity up","Days-weeks"),
        ("Stagnation",             "Slow chloramine decay, slight turbidity rise","Hours"),
        ("Operational Change",     "Sudden chloramine shift, turbidity stable","Minutes"),
        ("Sensor Fault",           "Zero variance detected on sensor readings","Variable"),
    ]

    # Header
    add_textbox(slide, 0.8, 1.9, 3.5, 0.35, "Cause Class", font_size=14, color=ACCENT, bold=True)
    add_textbox(slide, 4.5, 1.9, 4.5, 0.35, "Key Indicators", font_size=14, color=ACCENT, bold=True)
    add_textbox(slide, 9.2, 1.9, 2.5, 0.35, "Typical Duration", font_size=14, color=ACCENT, bold=True)

    for i, (cause, indicators, duration) in enumerate(causes):
        y = 2.4 + i * 0.42
        add_textbox(slide, 0.8, y, 3.5, 0.4, cause, font_size=13, color=WHITE, bold=True)
        add_textbox(slide, 4.5, y, 4.5, 0.4, indicators, font_size=12, color=LIGHT)
        add_textbox(slide, 9.2, y, 2.5, 0.4, duration, font_size=12, color=LIGHT)

    # Method explanation
    add_rounded_rect(slide, 0.5, 5.0, 3.8, 2.0)
    add_textbox(slide, 0.7, 5.1, 3.5, 0.4, "Rule Engine", font_size=16, color=ACCENT, bold=True)
    add_bullet_list(slide, 0.7, 5.5, 3.5, 1.4, [
        "6 domain-expert rules",
        "Always runs (no model needed)",
        "Provides human-readable",
        "  explanations for operators",
    ], font_size=12)

    add_rounded_rect(slide, 4.7, 5.0, 3.8, 2.0)
    add_textbox(slide, 4.9, 5.1, 3.5, 0.4, "ML Classifier", font_size=16, color=ACCENT, bold=True)
    add_bullet_list(slide, 4.9, 5.5, 3.5, 1.4, [
        "LightGBM 7-class model",
        "279 engineered features",
        "Probabilistic output",
        "  ranks all possible causes",
    ], font_size=12)

    add_rounded_rect(slide, 8.9, 5.0, 3.8, 2.0)
    add_textbox(slide, 9.1, 5.1, 3.5, 0.4, "SHAP Explanations", font_size=16, color=ACCENT, bold=True)
    add_bullet_list(slide, 9.1, 5.5, 3.5, 1.4, [
        "Per-prediction feature",
        "  importance via TreeSHAP",
        "Top-5 contributing factors",
        "  shown on dashboard",
    ], font_size=12)

    # ═══════════════════════════════════════════════════
    # SLIDE 6: Results & Event Scenarios
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 10, 0.7,
                "Results & Event Detection", font_size=36, color=ACCENT, bold=True)

    # Metrics cards
    metrics = [
        ("0.887", "AUC-ROC", "Potability\nClassifier"),
        ("0.86", "F1-macro", "Potability\nClassifier"),
        ("0.60", "Macro-F1", "Cause\nClassifier"),
        ("0.10", "MAE", "Chloramine\nForecaster"),
        ("~15 min", "Fastest", "Time-to-\nDetect"),
    ]
    for i, (val, metric, desc) in enumerate(metrics):
        x = 0.5 + i * 2.5
        add_rounded_rect(slide, x, 1.3, 2.2, 1.6)
        add_textbox(slide, x + 0.1, 1.4, 2.0, 0.6,
                    val, font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.0, 2.0, 0.35,
                    metric, font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.3, 2.0, 0.5,
                    desc, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Event scenarios
    add_textbox(slide, 0.8, 3.3, 10, 0.5,
                "3 Validated Event Scenarios", font_size=22, color=WHITE, bold=True)

    scenarios = [
        ("Scenario 1: Disinfection Failure",
         "Chloramine drops 4.0 to 0.5 mg/L over 2hrs.\n"
         "Detected at T+40min (WARNING), T+60min (CRITICAL).\n"
         "Cause: 'Disinfectant Decay' -- confidence 85%.\n"
         "Lead time: ~40 min before non-potable."),
        ("Scenario 2: Pipe Break Intrusion",
         "Turbidity jumps 2 to 6 NTU, conductivity spikes.\n"
         "Detected at T+15min (WARNING), T+20min (CRITICAL).\n"
         "Cause: 'Contamination Intrusion' -- confidence 80%.\n"
         "Lead time: ~15 min (fast onset)."),
        ("Scenario 3: Dead-End Stagnation",
         "Overnight chloramine decay in low-flow pipe.\n"
         "Detected at T+6h (CAUTION), T+7h (WARNING).\n"
         "Cause: 'Stagnation' -- confidence 70%.\n"
         "Lead time: ~2 hrs before violation."),
    ]
    for i, (title, body) in enumerate(scenarios):
        x = 0.5 + i * 4.2
        add_rounded_rect(slide, x, 3.9, 3.9, 3.0)
        add_textbox(slide, x + 0.15, 4.0, 3.6, 0.45,
                    title, font_size=14, color=ACCENT, bold=True)
        add_textbox(slide, x + 0.15, 4.5, 3.6, 2.3,
                    body, font_size=11, color=LIGHT)

    # ═══════════════════════════════════════════════════
    # SLIDE 7: Impact & Next Steps
    # ═══════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.4, 10, 0.7,
                "Impact & Next Steps", font_size=36, color=ACCENT, bold=True)

    # Impact section
    add_rounded_rect(slide, 0.5, 1.3, 5.8, 3.2)
    add_textbox(slide, 0.7, 1.4, 5.4, 0.5,
                "Public Health Impact", font_size=22, color=GREEN, bold=True)
    add_bullet_list(slide, 0.7, 2.0, 5.4, 2.4, [
        "15-40 min early warning before unsafe water reaches consumers",
        "~80% false alarm reduction vs. simple thresholds",
        "Automated root-cause identification saves hours of investigation",
        "5-site monitoring with 5-min granularity (scalable to 100+ sites)",
        "Interpretable explanations build operator trust",
        "Low-cost: runs on standard hardware, open-source stack",
    ], font_size=14)

    # Tech stack
    add_rounded_rect(slide, 6.8, 1.3, 5.8, 1.5)
    add_textbox(slide, 7.0, 1.4, 5.4, 0.4,
                "Technology Stack", font_size=18, color=ACCENT, bold=True)
    add_bullet_list(slide, 7.0, 1.8, 5.4, 0.9, [
        "Python | LightGBM | Scikit-learn | SHAP",
        "FastAPI (inference) | Streamlit (dashboard)",
        "Parquet (storage) | Plotly (visualization)",
    ], font_size=13)

    # Next steps
    add_rounded_rect(slide, 6.8, 3.1, 5.8, 3.4)
    add_textbox(slide, 7.0, 3.2, 5.4, 0.4,
                "Production Roadmap", font_size=18, color=ACCENT, bold=True)
    add_bullet_list(slide, 7.0, 3.7, 5.4, 2.7, [
        "Phase 1: Real sensor integration (MQTT / OPC-UA / SCADA)",
        "Phase 2: Weather & demand data fusion for better forecasting",
        "Phase 3: LSTM / Temporal Fusion Transformer for long horizons",
        "Phase 4: Graph Neural Network for pipe network topology",
        "Phase 5: Edge deployment (ONNX on IoT gateways)",
        "Phase 6: Docker + Kubernetes + Azure IoT Hub at scale",
    ], font_size=13)

    # Bottom bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, Inches(7.0), prs.slide_width, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_CARD
    bar.line.fill.background()
    add_textbox(slide, 1.0, 7.0, 11, 0.45,
                "AquaGuard  --  AI Night Challenge 2026  --  Protecting Public Health with Real-Time AI",
                font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # Save
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"PowerPoint saved: {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    build_presentation()
